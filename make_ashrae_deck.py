"""Generate a plain-language PowerPoint explaining the Red5-AHU solution,
how it satisfies ASHRAE 55 / 62.1 / 90.1 / Guideline 36 (Venn diagram +
comparison chart), and why it is an essential BMS tool.

Run:  python3 make_ashrae_deck.py
Out:  Red5-AHU-ASHRAE-Overview.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette -------------------------------------------------------------
INK      = RGBColor(0x0F, 0x17, 0x2A)
MUTED    = RGBColor(0x47, 0x55, 0x69)
FAINT    = RGBColor(0x64, 0x74, 0x8B)
ACCENT   = RGBColor(0x1D, 0x4E, 0xD8)   # blue  (Energy / brand)
GREEN    = RGBColor(0x04, 0x78, 0x57)   # fresh air
AMBER    = RGBColor(0xB4, 0x53, 0x09)   # comfort
PAGE     = RGBColor(0xF8, 0xFA, 0xFC)
CARD     = RGBColor(0xFF, 0xFF, 0xFF)
LINE     = RGBColor(0xE2, 0xE8, 0xF0)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
SLATE100 = RGBColor(0xF1, 0xF5, 0xF9)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)   # 16:9

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide(bg=PAGE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    # push background rectangle to back
    sp = r._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2, sp)
    return s


def box(s, l, t, w, h, text, size=18, bold=False, color=INK,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
        line_spacing=1.06):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        f = r.font
        f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = font
    return tb


def rect(s, l, t, w, h, fill=CARD, line=LINE, line_w=1.0, radius=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        l, t, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def set_alpha(shape, alpha_pct):
    """Set fill transparency (0..100 = fully transparent..opaque handled as
    alpha value). alpha_pct here = opacity percent."""
    sF = shape.fill._xPr.find(qn('a:solidFill'))
    srgb = sF.find(qn('a:srgbClr'))
    for old in srgb.findall(qn('a:alpha')):
        srgb.remove(old)
    a = srgb.makeelement(qn('a:alpha'), {'val': str(int(alpha_pct * 1000))})
    srgb.append(a)


def accent_header(s, kicker, title, tcolor=INK):
    rect(s, Inches(0), Inches(0), EMU_W, Inches(0.16), fill=ACCENT, line=None,
         radius=False)
    if kicker:
        box(s, Inches(0.7), Inches(0.42), Inches(11.9), Inches(0.4),
            kicker.upper(), size=13, bold=True, color=ACCENT)
    box(s, Inches(0.7), Inches(0.78), Inches(11.9), Inches(1.0),
        title, size=30, bold=True, color=tcolor)


def bullet_card(s, l, t, w, h, heading, hcolor, items, hsize=15, isize=13):
    rect(s, l, t, w, h, fill=CARD, line=LINE)
    box(s, l + Inches(0.22), t + Inches(0.18), w - Inches(0.44), Inches(0.5),
        heading, size=hsize, bold=True, color=hcolor)
    tb = s.shapes.add_textbox(l + Inches(0.22), t + Inches(0.72),
                              w - Inches(0.44), h - Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.06; p.space_after = Pt(6)
        r = p.add_run(); r.text = "•  " + it
        r.font.size = Pt(isize); r.font.color.rgb = MUTED; r.font.name = "Calibri"


# =========================================================================
# 1 — TITLE
# =========================================================================
s = slide(bg=RGBColor(0x0B, 0x12, 0x24))
rect(s, Inches(0), Inches(0), EMU_W, Inches(0.22), fill=ACCENT, line=None,
     radius=False)
box(s, Inches(0.9), Inches(2.05), Inches(11.6), Inches(1.4),
    "Red5-AHU", size=58, bold=True, color=WHITE)
box(s, Inches(0.9), Inches(3.25), Inches(11.6), Inches(1.0),
    "One chart to keep buildings comfortable, healthy, and energy-smart",
    size=25, bold=False, color=RGBColor(0x93, 0xC5, 0xFD))
box(s, Inches(0.9), Inches(4.35), Inches(11.6), Inches(1.2),
    "A plain-language look at how a smart air-handler tool meets every major "
    "ASHRAE standard — and why every Building Management System needs it.",
    size=16, color=RGBColor(0xCB, 0xD5, 0xE1))
for i, tg in enumerate(["Edge-native", "No cloud", "BACnet-native", "ASHRAE-aligned"]):
    p = rect(s, Inches(0.9 + i * 2.35), Inches(5.85), Inches(2.15), Inches(0.5),
             fill=RGBColor(0x1E, 0x29, 0x3B), line=RGBColor(0x33, 0x41, 0x55))
    box(s, Inches(0.9 + i * 2.35), Inches(5.9), Inches(2.15), Inches(0.4),
        tg, size=13, bold=True, color=RGBColor(0x93, 0xC5, 0xFD),
        align=PP_ALIGN.CENTER)

# =========================================================================
# 2 — THE PROBLEM (plain terms)
# =========================================================================
s = slide()
accent_header(s, "The challenge", "Three things every building must do — at once")
goals = [
    ("Feel comfortable", "Not too hot, not too cold, not too humid.", AMBER),
    ("Stay healthy", "Enough fresh outside air so the air is safe to breathe.", GREEN),
    ("Save energy", "Do both of the above without wasting power or money.", ACCENT),
]
cw = Inches(3.9); gap = Inches(0.25); x0 = Inches(0.7); y = Inches(2.15)
for i, (h, d, c) in enumerate(goals):
    l = x0 + i * (cw + gap)
    rect(s, l, y, cw, Inches(2.5), fill=CARD, line=LINE)
    rect(s, l, y, Inches(0.14), Inches(2.5), fill=c, line=None, radius=False)
    box(s, l + Inches(0.35), y + Inches(0.35), cw - Inches(0.6), Inches(0.6),
        h, size=20, bold=True, color=c)
    box(s, l + Inches(0.35), y + Inches(1.05), cw - Inches(0.6), Inches(1.3),
        d, size=15, color=MUTED)
rect(s, Inches(0.7), Inches(5.1), Inches(11.93), Inches(1.4),
     fill=RGBColor(0xFF, 0xFB, 0xEB), line=RGBColor(0xFD, 0xE6, 0x8A))
box(s, Inches(1.0), Inches(5.3), Inches(11.3), Inches(0.5),
    "The catch", size=15, bold=True, color=AMBER)
box(s, Inches(1.0), Inches(5.72), Inches(11.4), Inches(0.7),
    "These goals fight each other. More comfort and more fresh air usually cost "
    "more energy. Getting all three right, all day, in every season — is hard.",
    size=15, color=INK)

# =========================================================================
# 3 — FOUR ASHRAE RULEBOOKS (plain language)
# =========================================================================
s = slide()
accent_header(s, "The rulebooks", "Four ASHRAE standards, in one line each")
rows = [
    ("ASHRAE 55", "Comfort", "Will people feel comfortable?", AMBER),
    ("ASHRAE 62.1", "Fresh air", "Is there enough outside air to stay healthy?", GREEN),
    ("ASHRAE 90.1", "Energy", "Are we doing it without wasting energy?", ACCENT),
    ("Guideline 36", "The referee", "How do we run the equipment to hit all three?", INK),
]
y = Inches(2.15); rh = Inches(1.02); x0 = Inches(0.7); tot = Inches(11.93)
for i, (code, nick, q, c) in enumerate(rows):
    t = y + i * (rh + Inches(0.12))
    rect(s, x0, t, tot, rh, fill=CARD, line=LINE)
    rect(s, x0, t, Inches(2.5), rh, fill=c, line=None)
    box(s, x0 + Inches(0.2), t, Inches(2.1), rh, code, size=18, bold=True,
        color=WHITE, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    box(s, x0 + Inches(2.75), t, Inches(2.2), rh, nick, size=16, bold=True,
        color=c, anchor=MSO_ANCHOR.MIDDLE)
    box(s, x0 + Inches(5.0), t, tot - Inches(5.2), rh, q, size=15.5,
        color=MUTED, anchor=MSO_ANCHOR.MIDDLE)

# =========================================================================
# 4 — VENN DIAGRAM
# =========================================================================
s = slide()
accent_header(s, "The big idea", "One picture where all four rulebooks meet")
# three overlapping circles (center-based geometry so all fit the slide)
cxf, cyf = 4.35, 4.45          # venn center, inches
df = 2.95                      # diameter, inches
hx, hy = 1.02, 0.62            # center offsets, inches
D = Inches(df); half = df / 2.0
centers = [
    (cxf - hx, cyf - hy, AMBER),   # comfort  (top-left)
    (cxf + hx, cyf - hy, GREEN),   # fresh air(top-right)
    (cxf,      cyf + hy, ACCENT),  # energy   (bottom)
]
for ccx, ccy, c in centers:
    ov = s.shapes.add_shape(MSO_SHAPE.OVAL,
                            Inches(ccx - half), Inches(ccy - half), D, D)
    ov.fill.solid(); ov.fill.fore_color.rgb = c; set_alpha(ov, 40)
    ov.line.color.rgb = c; ov.line.width = Pt(1.5)
    ov.shadow.inherit = False
# outer labels
box(s, Inches(0.55), Inches(2.55), Inches(1.75), Inches(0.9),
    "COMFORT\nASHRAE 55", size=13, bold=True, color=AMBER,
    align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
box(s, Inches(6.55), Inches(2.55), Inches(1.9), Inches(0.9),
    "FRESH AIR\nASHRAE 62.1", size=13, bold=True, color=GREEN,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
box(s, Inches(2.85), Inches(6.45), Inches(3.0), Inches(0.7),
    "ENERGY · ASHRAE 90.1", size=13, bold=True, color=ACCENT,
    align=PP_ALIGN.CENTER)
# center label (the intersection)
box(s, Inches(cxf - 1.15), Inches(cyf - 0.5), Inches(2.3), Inches(1.0),
    "Psychrometric\nChart", size=16, bold=True, color=INK,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# right explainer
rx = Inches(9.3); rw = Inches(3.35)
rect(s, rx, Inches(2.0), rw, Inches(4.6), fill=CARD, line=LINE)
box(s, rx + Inches(0.25), Inches(2.2), rw - Inches(0.5), Inches(0.5),
    "Why it works", size=16, bold=True, color=ACCENT)
box(s, rx + Inches(0.25), Inches(2.75), rw - Inches(0.5), Inches(2.6),
    "All four standards live on the SAME chart — temperature across the "
    "bottom, humidity up the side.\n\n"
    "Plot one dot for the air and you can instantly see comfort, fresh-air "
    "mix, and energy — all together.",
    size=13.5, color=MUTED)
rect(s, rx + Inches(0.0), Inches(5.55), rw, Inches(1.05),
     fill=RGBColor(0xEC, 0xFD, 0xF5), line=RGBColor(0xA7, 0xF3, 0xD0))
box(s, rx + Inches(0.25), Inches(5.68), rw - Inches(0.5), Inches(0.9),
    "Guideline 36 is the referee — the control logic that keeps the dot in "
    "the sweet spot automatically.", size=12.5, bold=False, color=GREEN)

# =========================================================================
# 5 — HOW ONE CHART ANSWERS ALL FOUR
# =========================================================================
s = slide()
accent_header(s, "How it satisfies each", "One dot on the chart answers every question")
items = [
    ("55 — Comfort", "The comfort zone is drawn on the chart. Is the dot inside "
     "the good T/humidity band? If not, which way is it off?", AMBER),
    ("62.1 — Fresh air", "The mixed-air dot sits on the line between return and "
     "outside air — its position shows the fresh-air share at a glance.", GREEN),
    ("90.1 — Energy", "Enthalpy lines reveal free cooling: when outside air is "
     "cooler, the economizer meets the load with no mechanical cooling.", ACCENT),
    ("G36 — Sequences", "Visual proof the automatic sequences land the dot in "
     "the comfort zone, on the fresh-air line, at the lowest energy path.", INK),
]
x0 = Inches(0.7); cw = Inches(5.85); ch = Inches(1.9); gx = Inches(0.23); gy = Inches(0.22)
for i, (h, d, c) in enumerate(items):
    col = i % 2; row = i // 2
    l = x0 + col * (cw + gx); t = Inches(2.15) + row * (ch + gy)
    rect(s, l, t, cw, ch, fill=CARD, line=LINE)
    rect(s, l, t, Inches(0.13), ch, fill=c, line=None, radius=False)
    box(s, l + Inches(0.32), t + Inches(0.2), cw - Inches(0.55), Inches(0.5),
        h, size=17, bold=True, color=c)
    box(s, l + Inches(0.32), t + Inches(0.72), cw - Inches(0.6), ch - Inches(0.85),
        d, size=13.5, color=MUTED)
box(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.6),
    "In one sentence: plot where the air IS, see where it MUST be, how it gets "
    "there, and whether it does so at minimum energy.",
    size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# =========================================================================
# 6 — COMPARISON CHART
# =========================================================================
s = slide()
accent_header(s, "The difference", "Red5-AHU vs a traditional BMS")
# energy bars (normalized index) -- kept left of the feature table (x<6.6")
box(s, Inches(0.7), Inches(1.95), Inches(6.0), Inches(0.5),
    "Supply-air energy (lower is better)", size=15, bold=True, color=INK)
bx = 0.7; blabw = 2.35; bw100 = 3.4; bw28 = bw100 * 0.28
bar_left = Inches(bx + blabw); bh = Inches(0.62)
# traditional 100
box(s, Inches(bx), Inches(2.55), Inches(blabw), bh, "Fixed setpoint", size=13,
    color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
rect(s, bar_left, Inches(2.55), Inches(bw100), bh,
     fill=RGBColor(0x64, 0x74, 0x8B), line=None)
box(s, bar_left, Inches(2.55), Inches(bw100) - Inches(0.18), bh, "100",
    size=14, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
# red5 28
box(s, Inches(bx), Inches(3.35), Inches(blabw), bh, "Red5-AHU (B1–B10)", size=13,
    color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
rect(s, bar_left, Inches(3.35), Inches(bw28), bh, fill=GREEN, line=None)
box(s, bar_left + Inches(bw28) + Inches(0.08), Inches(3.35), Inches(0.6), bh, "28",
    size=14, bold=True, color=GREEN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
box(s, Inches(0.7), Inches(4.2), Inches(6.0), Inches(1.0),
    "≈ 70–75% less supply-air energy than a fixed year-round setpoint "
    "(internal estimate, temperate climate).", size=12.5, color=FAINT)

# feature check table on right
fx = Inches(7.05); fw = Inches(5.6)
feats = [
    ("Where it runs", "Cloud / head-end PC", "On the controller (edge)"),
    ("Standards built in", "Add-on / manual", "ASHRAE 55 · 62.1 · 90.1 · G36"),
    ("Diagnose an AHU", "Dig through trends", "One glance at the chart"),
    ("Vendor lock-in", "Usually yes", "Open — MQTT/Modbus/BACnet"),
    ("Languages", "Often one", "Five"),
]
rect(s, fx, Inches(1.95), fw, Inches(4.7), fill=CARD, line=LINE)
# header
box(s, fx + Inches(0.2), Inches(2.08), Inches(1.9), Inches(0.5), "", size=11)
box(s, fx + Inches(2.05), Inches(2.08), Inches(1.7), Inches(0.5),
    "Traditional BMS", size=11.5, bold=True, color=FAINT, align=PP_ALIGN.CENTER)
box(s, fx + Inches(3.8), Inches(2.08), Inches(1.7), Inches(0.5),
    "Red5-AHU", size=11.5, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
ry = Inches(2.55)
for i, (k, a, b) in enumerate(feats):
    t = ry + i * Inches(0.82)
    if i % 2 == 0:
        rect(s, fx + Inches(0.08), t, fw - Inches(0.16), Inches(0.78),
             fill=RGBColor(0xFA, 0xFC, 0xFF), line=None)
    box(s, fx + Inches(0.2), t, Inches(1.85), Inches(0.78), k, size=12,
        bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    box(s, fx + Inches(2.02), t, Inches(1.75), Inches(0.78), a, size=11,
        color=MUTED, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    box(s, fx + Inches(3.78), t, Inches(1.72), Inches(0.78), b, size=11,
        bold=True, color=GREEN, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

# =========================================================================
# 6b — REAL RESULT (product chart image, if present)
# =========================================================================
_img = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                    "frontend", "public", "band_monthly_energy.png")
if os.path.exists(_img):
    s = slide()
    accent_header(s, "Seeing it in action", "The control follows the seasons")
    s.shapes.add_picture(_img, Inches(0.7), Inches(2.05), height=Inches(4.7))
    box(s, Inches(6.4), Inches(2.25), Inches(6.2), Inches(0.6),
        "What you're looking at", size=16, bold=True, color=ACCENT)
    box(s, Inches(6.4), Inches(2.85), Inches(6.25), Inches(3.6),
        "Instead of one fixed target all year, Red5-AHU sets the supply-air "
        "target to match the season (climate bands B1–B10).\n\n"
        "In mild months it leans on free outside-air cooling; in harsh months "
        "it does only as much conditioning as needed.\n\n"
        "The result: month after month the air handler uses far less energy "
        "than a fixed setpoint — automatically, with no operator effort.",
        size=15, color=MUTED)

# =========================================================================
# 7 — WHY ESSENTIAL FOR BMS
# =========================================================================
s = slide()
accent_header(s, "The bottom line", "Why it's an essential BMS tool")
cards = [
    ("Sees everything at once", "Comfort, fresh air, and energy on one live "
     "picture — no guesswork, no digging through data.", ACCENT),
    ("Saves real energy", "Automatic free-cooling and climate-adaptive supply "
     "air cut supply-air energy ~70%.", GREEN),
    ("Runs on the controller", "No cloud, no extra server. Fewer parts to fail, "
     "keeps working if the network drops.", AMBER),
    ("Speaks BMS natively", "BACnet-native and open (MQTT/Modbus) — no vendor "
     "lock-in, integrates with what you have.", ACCENT),
    ("Meets the standards", "Aligned to ASHRAE 55 / 62.1 / 90.1 and Guideline "
     "36 out of the box.", GREEN),
    ("Easy to roll out", "Diagnose at a glance, clone the setup to the next "
     "unit, in five languages.", AMBER),
]
x0 = Inches(0.7); cw = Inches(3.9); ch = Inches(1.95); gx = Inches(0.23); gy = Inches(0.22)
for i, (h, d, c) in enumerate(cards):
    col = i % 3; row = i // 3
    l = x0 + col * (cw + gx); t = Inches(2.1) + row * (ch + gy)
    rect(s, l, t, cw, ch, fill=CARD, line=LINE)
    rect(s, l, t, cw, Inches(0.1), fill=c, line=None, radius=False)
    box(s, l + Inches(0.28), t + Inches(0.28), cw - Inches(0.5), Inches(0.5),
        h, size=16, bold=True, color=c)
    box(s, l + Inches(0.28), t + Inches(0.85), cw - Inches(0.55), ch - Inches(1.0),
        d, size=13, color=MUTED)

# =========================================================================
# 8 — CLOSING
# =========================================================================
s = slide(bg=RGBColor(0x0B, 0x12, 0x24))
rect(s, Inches(0), Inches(0), EMU_W, Inches(0.22), fill=ACCENT, line=None,
     radius=False)
box(s, Inches(1.0), Inches(2.5), Inches(11.3), Inches(1.6),
    "Four rulebooks. One chart. Zero cloud.",
    size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, Inches(1.4), Inches(4.2), Inches(10.5), Inches(1.4),
    "Red5-AHU turns every air handler into a self-diagnosing, self-optimizing "
    "asset that satisfies ASHRAE 55, 62.1, 90.1 and Guideline 36 — right on "
    "the controller.", size=18, color=RGBColor(0xCB, 0xD5, 0xE1),
    align=PP_ALIGN.CENTER)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Red5-AHU-ASHRAE-Overview.pptx")
prs.save(out)
print("Wrote", out, "-", len(prs.slides._sldIdLst), "slides")
