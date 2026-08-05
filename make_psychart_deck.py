"""Generate a vendor-neutral PowerPoint about the psychrometric chart:
who uses it and why (HVAC equipment designers), why traditional BMS never
adopted it, and how the one chart answers ASHRAE 55 / 62.1 / 90.1 and
Guideline 36.

Run:  python3 make_psychart_deck.py
Out:  Psychart-HVAC-ASHRAE-Overview.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

try:
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
except ImportError:                                    # older python-pptx
    MSO_LINE_DASH_STYLE = None

# ---- palette -------------------------------------------------------------
INK      = RGBColor(0x0F, 0x17, 0x2A)
MUTED    = RGBColor(0x47, 0x55, 0x69)
FAINT    = RGBColor(0x64, 0x74, 0x8B)
ACCENT   = RGBColor(0x1D, 0x4E, 0xD8)   # blue  (energy / structure)
GREEN    = RGBColor(0x04, 0x78, 0x57)   # fresh air
AMBER    = RGBColor(0xB4, 0x53, 0x09)   # comfort
VIOLET   = RGBColor(0x7C, 0x3A, 0xED)   # equal-energy (enthalpy) lines
ROSE     = RGBColor(0xE1, 0x1D, 0x48)   # dry-bulb "open" but costlier
SLATE    = RGBColor(0x33, 0x41, 0x55)
PAGE     = RGBColor(0xF8, 0xFA, 0xFC)
CARD     = RGBColor(0xFF, 0xFF, 0xFF)
LINE     = RGBColor(0xE2, 0xE8, 0xF0)
GRID     = RGBColor(0xEE, 0xF2, 0xF7)
AXIS     = RGBColor(0x94, 0xA3, 0xB8)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)   # 16:9

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide(bg=PAGE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    sp = r._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2, sp)
    return s


def box(s, l, t, w, h, text, size=18, bold=False, color=INK,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
        line_spacing=1.06):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        f = r.font
        f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = font
    return tb


def rect(s, l, t, w, h, fill=CARD, line=LINE, line_w=1.0, radius=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        l, t, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def set_alpha(shape, opacity_pct):
    sF = shape.fill._xPr.find(qn('a:solidFill'))
    srgb = sF.find(qn('a:srgbClr'))
    for old in srgb.findall(qn('a:alpha')):
        srgb.remove(old)
    srgb.append(srgb.makeelement(qn('a:alpha'),
                                 {'val': str(int(opacity_pct * 1000))}))


def conn(s, x1, y1, x2, y2, color=INK, width=1.75, dash=False):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                               Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dash and MSO_LINE_DASH_STYLE is not None:
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return c


def dot(s, x, y, color, r=0.07):
    ov = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - r), Inches(y - r),
                            Inches(2 * r), Inches(2 * r))
    ov.fill.solid(); ov.fill.fore_color.rgb = color
    ov.line.color.rgb = WHITE; ov.line.width = Pt(1)
    ov.shadow.inherit = False
    return ov


# saturation curve, from ASHRAE saturation pressures at 101.325 kPa
# (dry-bulb degC, humidity ratio g/kg dry air)
SAT = [(0, 3.775), (5, 5.402), (10, 7.630), (15, 10.647), (20, 14.697),
       (25, 20.082), (30, 27.202), (31.5, 30.0)]


def wsat(T):
    if T >= SAT[-1][0]:
        return 99.0
    for (t1, w1), (t2, w2) in zip(SAT, SAT[1:]):
        if t1 <= T <= t2:
            return w1 + (w2 - w1) * (T - t1) / (t2 - t1)
    return 99.0


def enth_seg(T0, W0, k=0.402):
    """Endpoints of the constant-enthalpy line through (T0, W0), clipped to
    the saturation curve and the plot box. k = dW/dT at constant enthalpy."""
    pts = []
    T = 0.0
    while T <= 40.0001:
        W = W0 + k * (T0 - T)
        if 0.0 <= W <= 30.0 and W <= wsat(T):
            pts.append((T, W))
        T += 0.1
    return (pts[0], pts[-1]) if pts else None


def accent_header(s, kicker, title, tcolor=INK):
    rect(s, Inches(0), Inches(0), EMU_W, Inches(0.16), fill=ACCENT, line=None,
         radius=False)
    if kicker:
        box(s, Inches(0.7), Inches(0.42), Inches(11.9), Inches(0.4),
            kicker.upper(), size=13, bold=True, color=ACCENT)
    box(s, Inches(0.7), Inches(0.78), Inches(11.9), Inches(1.0),
        title, size=30, bold=True, color=tcolor)


def lead(s, text, top=1.55):
    box(s, Inches(0.7), Inches(top), Inches(11.9), Inches(0.75), text,
        size=15, color=MUTED, line_spacing=1.2)


def bullet_card(s, l, t, w, h, heading, hcolor, items, hsize=15, isize=12.5,
                intro=None):
    rect(s, l, t, w, h, fill=CARD, line=LINE)
    rect(s, l, t, Inches(0.14), h, fill=hcolor, line=None, radius=False)
    box(s, l + Inches(0.34), t + Inches(0.22), w - Inches(0.55), Inches(0.5),
        heading, size=hsize, bold=True, color=hcolor)
    ty = t + Inches(0.78)
    if intro:
        box(s, l + Inches(0.34), ty, w - Inches(0.6), Inches(0.6), intro,
            size=isize, color=MUTED, line_spacing=1.15)
        ty = ty + Inches(0.55)
    tb = s.shapes.add_textbox(l + Inches(0.34), ty, w - Inches(0.6),
                              h - (ty - t) - Inches(0.15))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.12; p.space_after = Pt(5)
        r = p.add_run(); r.text = "•  " + it
        r.font.size = Pt(isize); r.font.color.rgb = MUTED; r.font.name = "Calibri"


# =========================================================================
# 1 — TITLE
# =========================================================================
s = slide(bg=RGBColor(0x0B, 0x12, 0x24))
rect(s, Inches(0), Inches(0), EMU_W, Inches(0.22), fill=ACCENT, line=None,
     radius=False)
box(s, Inches(0.9), Inches(1.85), Inches(11.6), Inches(1.5),
    "The Psychrometric Chart", size=56, bold=True, color=WHITE)
box(s, Inches(0.9), Inches(3.35), Inches(11.6), Inches(1.0),
    "HVAC's master diagram — and the missing tool in most Building "
    "Management Systems",
    size=24, color=RGBColor(0x93, 0xC5, 0xFD), line_spacing=1.2)
box(s, Inches(0.9), Inches(4.65), Inches(11.4), Inches(1.2),
    "Why equipment designers cannot work without it, why traditional BMS "
    "never adopted it, and how this one chart answers ASHRAE 55, 62.1, 90.1 "
    "and Guideline 36 at the same time.",
    size=16, color=RGBColor(0xCB, 0xD5, 0xE1), line_spacing=1.25)
for i, tg in enumerate(["Moist-air states", "Design tool", "ASHRAE-aligned",
                        "Live controls"]):
    rect(s, Inches(0.9 + i * 2.55), Inches(6.15), Inches(2.35), Inches(0.5),
         fill=RGBColor(0x1E, 0x29, 0x3B), line=RGBColor(0x33, 0x41, 0x55))
    box(s, Inches(0.9 + i * 2.55), Inches(6.2), Inches(2.35), Inches(0.4),
        tg, size=13, bold=True, color=RGBColor(0x93, 0xC5, 0xFD),
        align=PP_ALIGN.CENTER)

# =========================================================================
# 2 — WHO USES IT, ON WHAT EQUIPMENT
# =========================================================================
s = slide()
accent_header(s, "Start here", "Who uses this chart — and on what equipment")
lead(s, "Long before a building is controlled, someone has to choose the "
        "machines that condition its air. Those people live on the "
        "psychrometric chart.")
cw = Inches(3.9); gap = Inches(0.25); x0 = Inches(0.7)
y = Inches(2.4); ch = Inches(3.5)
bullet_card(s, x0, y, cw, ch, "Who they are", ACCENT,
            ["manufacturers' application engineers",
             "consulting / design engineers",
             "commissioning engineers who prove it works"],
            intro="\u201cHVAC equipment designers\u201d are the mechanical "
                  "engineers who select and size air-conditioning equipment:")
bullet_card(s, x0 + cw + gap, y, cw, ch, "What \u201cHVAC equipment\u201d means",
            GREEN,
            ["air-handling units (AHUs) and their heating / cooling coils",
             "mixing & economizer dampers",
             "humidifiers and dehumidifiers",
             "energy-recovery wheels (ERV / HRV)",
             "chillers and DX systems, fans, VAV terminals"],
            intro="The machines that actually change the air:")
bullet_card(s, x0 + 2 * (cw + gap), y, cw, ch, "What they must decide", AMBER,
            ["How much cooling or heating is needed",
             "How much of that cooling removes MOISTURE rather than heat",
             "How big the coil must be",
             "How much outside air to mix in",
             "Whether condensation or coil freezing will occur"])
rect(s, Inches(0.7), Inches(6.25), Inches(11.93), Inches(0.8),
     fill=RGBColor(0xEF, 0xF6, 0xFF), line=RGBColor(0xBF, 0xDB, 0xFE))
box(s, Inches(0.98), Inches(6.35), Inches(11.4), Inches(0.6),
    "Every one of those decisions is a point or a line on the psychrometric "
    "chart. That is why the chart — not the equipment catalogue — is where "
    "air-conditioning design actually happens.",
    size=13.5, color=INK, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)

# =========================================================================
# 3 — WHY DESIGNERS RELY ON IT
# =========================================================================
s = slide()
accent_header(s, "Why it exists", "Why designers cannot work without it")
lead(s, "Air conditioning is not simply \u201cmaking air colder\u201d. It is "
        "moving MOIST AIR from one state to another — and the psychrometric "
        "chart is the map of every state air can be in.")
items = [
    ("Every process is a line",
     "Cooling, heating, humidifying, dehumidifying and mixing each move the "
     "air in a characteristic direction. Draw the line and you have described "
     "the process.", ACCENT),
    ("Capacity comes from the chart",
     "The energy content (enthalpy) at each end of the line, multiplied by "
     "airflow, is the kW the coil must deliver. That is how equipment gets "
     "sized.", GREEN),
    ("Sensible vs latent",
     "The chart separates cooling that lowers temperature from cooling that "
     "removes water. The two behave differently and must be sized "
     "separately.", AMBER),
    ("Physical limits are visible",
     "The curved boundary is 100% humidity — air cannot go beyond it. It "
     "shows where condensation begins and where a coil would freeze.", INK),
]
x0 = Inches(0.7); cw = Inches(5.85); ch = Inches(1.95)
gx = Inches(0.23); gy = Inches(0.22)
for i, (h, d, c) in enumerate(items):
    l = x0 + (i % 2) * (cw + gx); t = Inches(2.4) + (i // 2) * (ch + gy)
    rect(s, l, t, cw, ch, fill=CARD, line=LINE)
    rect(s, l, t, Inches(0.13), ch, fill=c, line=None, radius=False)
    box(s, l + Inches(0.32), t + Inches(0.22), cw - Inches(0.55), Inches(0.5),
        h, size=17, bold=True, color=c)
    box(s, l + Inches(0.32), t + Inches(0.74), cw - Inches(0.6),
        ch - Inches(0.88), d, size=13.5, color=MUTED, line_spacing=1.18)

# =========================================================================
# 4 — HOW TO READ IT (chart drawn from real saturation data)
# =========================================================================
s = slide()
accent_header(s, "How to read it", "One dot is the state of the air")

# plot mapping: x = dry-bulb 0..40 degC, y = humidity ratio 0..30 g/kg
X0, Y0 = 0.95, 5.55        # origin (inches)
SX, SY = 0.115, 0.105      # inches per degC / per g/kg
PX = lambda T: X0 + T * SX
PY = lambda W: Y0 - W * SY

rect(s, Inches(0.62), Inches(1.95), Inches(5.55), Inches(4.0),
     fill=CARD, line=LINE)

# faint grid
for T in (5, 10, 15, 20, 25, 30, 35):
    conn(s, PX(T), PY(30), PX(T), PY(0), color=GRID, width=0.75)
for W in (5, 10, 15, 20, 25):
    conn(s, PX(0), PY(W), PX(40), PY(W), color=GRID, width=0.75)
# axes
conn(s, PX(0), PY(0), PX(40), PY(0), color=AXIS, width=1.25)
conn(s, PX(0), PY(0), PX(0), PY(30), color=AXIS, width=1.25)

# lines of equal total energy (constant enthalpy), clipped at saturation
for (T0, W0), lw in [((13, 8.8), 0.9), ((33, 18.0), 0.9), ((24, 9.3), 1.4)]:
    sg = enth_seg(T0, W0)
    if sg:
        (t1, w1), (t2, w2) = sg
        conn(s, PX(t1), PY(w1), PX(t2), PY(w2), color=VIOLET, width=lw,
             dash=True)
box(s, Inches(PX(16.4) - 1.5), Inches(PY(13.06) - 0.12), Inches(1.5),
    Inches(0.25), "equal energy (enthalpy)", size=9, bold=True, color=VIOLET,
    align=PP_ALIGN.RIGHT)

# saturation curve (100% RH)
verts = [(Inches(PX(t)), Inches(PY(w))) for t, w in SAT]
ff = s.shapes.build_freeform(verts[0][0], verts[0][1], 1.0)
ff.add_line_segments(verts[1:], close=False)
curve = ff.convert_to_shape()
curve.fill.background()
curve.line.color.rgb = ACCENT
curve.line.width = Pt(2.25)
curve.shadow.inherit = False
box(s, Inches(PX(32)), Inches(PY(30) - 0.02), Inches(1.2), Inches(0.5),
    "100% RH\n(saturation)", size=9.5, bold=True, color=ACCENT)

# comfort zone (ASHRAE 55, approx 20-26 degC / 6-12 g/kg)
cz = rect(s, Inches(PX(20)), Inches(PY(12)), Inches(PX(26) - PX(20)),
          Inches(PY(6) - PY(12)), fill=AMBER, line=AMBER, line_w=1.25)
set_alpha(cz, 16)
box(s, Inches((PX(20) + PX(26)) / 2 - 1.3), Inches(PY(6) + 0.05), Inches(2.6),
    Inches(0.3), "Comfort zone (ASHRAE 55)", size=9.5, bold=True, color=AMBER,
    align=PP_ALIGN.CENTER)

# processes: OA-RA mixing line, then coil MA -> SA
OA, RA = (33, 18.0), (24, 9.3)
MA, SA = (26.7, 11.91), (13, 8.8)
conn(s, PX(OA[0]), PY(OA[1]), PX(RA[0]), PY(RA[1]), color=GREEN, width=1.75,
     dash=True)
conn(s, PX(MA[0]), PY(MA[1]), PX(SA[0]), PY(SA[1]), color=ACCENT, width=2.25)
for (T, W), c, nm, dx, dy in [
        (OA, AMBER, "OA", 0.10, -0.17),
        (RA, GREEN, "RA", 0.09, 0.03),
        (MA, INK, "MA", 0.09, -0.20),
        (SA, ACCENT, "SA", -0.40, -0.20)]:
    dot(s, PX(T), PY(W), c, r=0.06)
    box(s, Inches(PX(T) + dx), Inches(PY(W) + dy), Inches(0.45), Inches(0.25),
        nm, size=10.5, bold=True, color=c)
# axis captions
box(s, Inches(PX(0)), Inches(PY(0) + 0.10), Inches(PX(40) - PX(0)),
    Inches(0.28), "Dry-bulb temperature  \u2192", size=10.5, bold=True,
    color=MUTED, align=PP_ALIGN.CENTER)
box(s, Inches(0.75), Inches(PY(30) - 0.26), Inches(2.6), Inches(0.28),
    "\u2191  Moisture in the air", size=10.5, bold=True, color=MUTED)
box(s, Inches(0.62), Inches(6.02), Inches(5.6), Inches(1.05),
    "OA outside air  ·  RA return air  ·  MA the mixture entering the coil  ·  "
    "SA supply air delivered to the rooms. Dashed green = mixing; solid blue = "
    "the cooling coil. The violet diagonals are lines of equal total energy: "
    "outside air here holds 79 kJ/kg against return air's 48, so it sits well "
    "above RA's line and costs more to condition — no free cooling today.",
    size=10.5, color=FAINT, line_spacing=1.2)

# reading guide — 2 x 3 cards on the right
guide = [
    ("Two axes fix everything",
     "Across is temperature; up is the actual weight of water the air "
     "carries. Pin both and the air's state is pinned — relative humidity, "
     "dew point, wet-bulb and energy content all follow from that one dot.",
     ACCENT),
    ("The curve is a hard limit",
     "It marks 100% humidity — air can never sit above it. Drive the dot onto "
     "the curve and water comes out: the dew point, a wet coil, a fogged "
     "window.", ACCENT),
    ("Mixing is a straight line",
     "Outside and return air blend to a point on the line joining them. How "
     "far along it sits gives the proportion: a third of the way from RA "
     "means a third outside air.", GREEN),
    ("A coil pulls down and left",
     "Left removes temperature (sensible heat); down removes water (latent "
     "heat). The slope of MA→SA is the split — and that split is what sizing "
     "a coil actually means.", ACCENT),
    ("Comfort is an area",
     "The acceptable range of temperature and humidity is drawn straight onto "
     "the chart. Supply air must be chosen so the room lands inside that box, "
     "not merely at the right temperature.", AMBER),
    ("Faults show as bad geometry",
     "If the mixed-air dot is not on the line between outside and return air, "
     "something is lying — a stuck damper or a drifting sensor. The picture "
     "catches what alarm limits miss.", INK),
]
gx0 = 6.35; gcw = 2.98; gch = 1.6; ggx = 0.16; ggy = 0.14
for i, (h, d, c) in enumerate(guide):
    l = Inches(gx0 + (i % 2) * (gcw + ggx))
    t = Inches(2.0 + (i // 2) * (gch + ggy))
    rect(s, l, t, Inches(gcw), Inches(gch), fill=CARD, line=LINE)
    ov = s.shapes.add_shape(MSO_SHAPE.OVAL, l + Inches(0.18), t + Inches(0.16),
                            Inches(0.24), Inches(0.24))
    ov.fill.solid(); ov.fill.fore_color.rgb = c
    ov.line.fill.background(); ov.shadow.inherit = False
    box(s, l + Inches(0.18), t + Inches(0.16), Inches(0.24), Inches(0.24),
        str(i + 1), size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)
    box(s, l + Inches(0.5), t + Inches(0.15), Inches(gcw - 0.65), Inches(0.3),
        h, size=12.5, bold=True, color=c)
    box(s, l + Inches(0.2), t + Inches(0.5), Inches(gcw - 0.4),
        Inches(gch - 0.6), d, size=10, color=MUTED, line_spacing=1.12)

# =========================================================================
# 5 — WORKED EXAMPLE: WHEN IS OUTSIDE AIR FREE?
# =========================================================================
s = slide()
accent_header(s, "Worked example", "When is outside air free?")

EX0, EY0 = 0.95, 5.50
ESX, ESY = 0.115, 0.100
ex = lambda T: EX0 + T * ESX
ey = lambda W: EY0 - W * ESY

rect(s, Inches(0.62), Inches(1.95), Inches(5.55), Inches(3.95), fill=CARD,
     line=LINE)


def poly(s, pts, fill, opacity):
    v = [(Inches(ex(T)), Inches(ey(W))) for T, W in pts]
    b = s.shapes.build_freeform(v[0][0], v[0][1], 1.0)
    b.add_line_segments(v[1:], close=True)
    shp = b.convert_to_shape()
    shp.fill.solid(); shp.fill.fore_color.rgb = fill; set_alpha(shp, opacity)
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


# the four decision regions, split by the RA enthalpy line and the RA dry-bulb
poly(s, [(16.9, 12.15), (15, 10.647), (10, 7.630), (5, 5.402), (0, 3.775),
         (0, 0), (24, 0), (24, 9.3)], GREEN, 13)
poly(s, [(16.9, 12.15), (20, 14.697), (24, 19.005), (24, 9.3)], ROSE, 16)
poly(s, [(24, 9.3), (24, 0), (40, 0), (40, 2.87)], AMBER, 16)
poly(s, [(24, 9.3), (24, 19.005), (25, 20.082), (30, 27.202), (31.5, 30),
         (40, 30), (40, 2.87)], FAINT, 10)

conn(s, ex(0), ey(0), ex(40), ey(0), color=AXIS, width=1.25)
conn(s, ex(0), ey(0), ex(0), ey(30), color=AXIS, width=1.25)
# the two competing changeover lines
conn(s, ex(24), ey(19.005), ex(24), ey(0), color=SLATE, width=1.5, dash=True)
conn(s, ex(16.9), ey(12.15), ex(40), ey(2.87), color=VIOLET, width=2.0)
# saturation
_v = [(Inches(ex(t)), Inches(ey(w))) for t, w in SAT]
_b = s.shapes.build_freeform(_v[0][0], _v[0][1], 1.0)
_b.add_line_segments(_v[1:], close=False)
_c = _b.convert_to_shape()
_c.fill.background(); _c.line.color.rgb = ACCENT; _c.line.width = Pt(2.25)
_c.shadow.inherit = False


def lab(T, W, text, size=10, color=INK, w=1.7, bold=True):
    box(s, Inches(ex(T) - 0.04), Inches(ey(W) - 0.13), Inches(w), Inches(0.26),
        text, size=size, bold=bold, color=color)


for (T, W), c, nm, r in [((24, 9.3), GREEN, "RA 48", 0.07),
                         ((22, 14.96), ROSE, "A 60", 0.07),
                         ((28, 2.33), AMBER, "B 34", 0.07),
                         ((33, 18.0), FAINT, "OA 79", 0.055)]:
    dot(s, ex(T), ey(W), c, r=r)
lab(24.5, 7.4, "RA 48", 11, GREEN, w=0.9)
lab(18.3, 15.6, "A 60", 11, ROSE, w=0.8)
lab(28.8, 1.9, "B 34", 11, AMBER, w=0.8)
lab(33.8, 18.4, "OA 79", 10, FAINT, w=0.9)
lab(5.6, 1.9, "FREE COOLING", 10.5, GREEN, w=1.6)
lab(5.6, 0.55, "both methods agree", 9, GREEN, w=1.7, bold=False)
lab(31.4, 23.9, "hot & humid —", 9.5, FAINT, w=1.3)
lab(31.4, 22.5, "minimum OA", 9.5, FAINT, w=1.3)
lab(6.6, 13.2, "equal energy as RA", 9.5, VIOLET, w=1.6)
lab(24.5, 16.7, "dry-bulb only", 9.5, SLATE, w=1.3)

box(s, Inches(0.62), Inches(5.95), Inches(5.6), Inches(0.5),
    "Numbers are total energy in kJ/kg.  RA 24 °C / 50% RH = 48  ·  "
    "A 22 °C / 90% RH = 60  ·  B 28 °C / 10% RH = 34  ·  OA 33 °C / 18 g/kg "
    "= 79 (the outside air from the previous slide).",
    size=10, color=FAINT, line_spacing=1.2)

# right-hand notes
nx = Inches(6.4); nw = Inches(6.23)
box(s, nx, Inches(2.0), nw, Inches(0.4), "Two lines, four outcomes",
    size=16, bold=True, color=ACCENT)
outcomes = [
    ("1", GREEN, "Below-left of the violet line — outside air holds less "
     "energy than the air coming back. Open the economizer and let mechanical "
     "cooling back off."),
    ("2", FAINT, "Above-right — outside air is a load, not a resource. Close "
     "to the minimum ventilation rate 62.1 demands and let the coil do the "
     "work."),
    ("A", ROSE, "Cool but muggy. At 22 °C it is 2° cooler than the return, so "
     "a dry-bulb economizer opens — yet it carries 60 kJ/kg against 48. You "
     "have just imported latent load for the coil to wring back out."),
    ("B", AMBER, "Warm but dry. At 28 °C dry-bulb keeps it shut, yet 34 kJ/kg "
     "is far below the return. Energy says use it — though a dry-bulb high "
     "limit still guards this corner when the load is mostly sensible."),
    ("=", VIOLET, "On the line the two streams are energy-neutral, so the "
     "choice falls to humidity and fan power. Sequences add a deadband here "
     "so the dampers do not hunt."),
]
ny = 2.5
for mark, c, txt in outcomes:
    ov = s.shapes.add_shape(MSO_SHAPE.OVAL, nx, Inches(ny), Inches(0.26),
                            Inches(0.26))
    ov.fill.solid(); ov.fill.fore_color.rgb = c
    ov.line.fill.background(); ov.shadow.inherit = False
    box(s, nx, Inches(ny), Inches(0.26), Inches(0.26), mark, size=10.5,
        bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)
    box(s, nx + Inches(0.4), Inches(ny - 0.05), nw - Inches(0.42),
        Inches(0.72), txt, size=11.5, color=MUTED, line_spacing=1.16)
    ny += 0.76

rect(s, Inches(0.7), Inches(6.5), Inches(11.93), Inches(0.8),
     fill=RGBColor(0xEF, 0xF6, 0xFF), line=RGBColor(0xBF, 0xDB, 0xFE))
box(s, Inches(0.98), Inches(6.6), Inches(11.4), Inches(0.6),
    "A dry-bulb sensor can only draw the VERTICAL line; the chart draws the "
    "DIAGONAL. The two shaded wedges are precisely where they disagree — and "
    "where energy is quietly lost.",
    size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)

# =========================================================================
# 6 — WHY TRADITIONAL BMS NEVER ADOPTED IT
# =========================================================================
s = slide()
accent_header(s, "The gap", "Why traditional BMS never adopted the chart")
reasons = [
    ("Controls think in single numbers",
     "A BMS is assembled from loops that each watch ONE value — a "
     "temperature, a pressure. A state point needs two readings judged "
     "together; classic controls had no concept of a two-dimensional air "
     "state.", ACCENT),
    ("The sensors often aren't there",
     "You need temperature AND humidity at the same place. Many units "
     "measure only dry-bulb on outside, mixed and return air — extra humidity "
     "sensors mean extra cost, wiring and calibration drift.", GREEN),
    ("The maths was too expensive",
     "Moisture content and energy content are non-linear and depend on "
     "barometric pressure. Older controllers had very little memory and no "
     "floating-point maths, so psychrometrics were left out.", AMBER),
    ("Graphics couldn't draw it",
     "Early front-ends drew simple equipment schematics and time-series "
     "trends. A live chart with curves and overlays needs real rendering — so "
     "the chart stayed a static page in the design binder.", ACCENT),
    ("Design and operations are separate worlds",
     "The chart is used once, offline, to select equipment. What is handed to "
     "the controls contractor is a list of setpoints — not the physics behind "
     "them. The reasoning is lost at handover.", GREEN),
    ("The simpler option was allowed",
     "Codes permit dry-bulb-only economizer switching, so designs avoided "
     "humidity sensing that was historically unreliable. The easier, "
     "compliant path quietly became the default.", AMBER),
]
x0 = Inches(0.7); cw = Inches(3.9); ch = Inches(2.05)
gx = Inches(0.2); gy = Inches(0.2)
for i, (h, d, c) in enumerate(reasons):
    l = x0 + (i % 3) * (cw + gx); t = Inches(2.0) + (i // 3) * (ch + gy)
    rect(s, l, t, cw, ch, fill=CARD, line=LINE)
    rect(s, l, t, cw, Inches(0.1), fill=c, line=None, radius=False)
    box(s, l + Inches(0.26), t + Inches(0.22), cw - Inches(0.5), Inches(0.55),
        h, size=14.5, bold=True, color=c)
    box(s, l + Inches(0.26), t + Inches(0.82), cw - Inches(0.52),
        ch - Inches(0.92), d, size=11.5, color=MUTED, line_spacing=1.14)
rect(s, Inches(0.7), Inches(6.35), Inches(11.93), Inches(0.85),
     fill=RGBColor(0xFF, 0xFB, 0xEB), line=RGBColor(0xFD, 0xE6, 0x8A))
box(s, Inches(0.98), Inches(6.45), Inches(11.4), Inches(0.3),
    "The consequence", size=13, bold=True, color=AMBER)
box(s, Inches(0.98), Inches(6.74), Inches(11.5), Inches(0.4),
    "Operators end up diagnosing a TWO-DIMENSIONAL problem — heat AND "
    "moisture — by staring at separate one-dimensional trend lines. The "
    "information is there; the picture is not.", size=13, color=INK)

# =========================================================================
# 6 — FOUR ASHRAE RULEBOOKS
# =========================================================================
s = slide()
accent_header(s, "The rulebooks", "Four ASHRAE standards, in one line each")
rows = [
    ("ASHRAE 55", "Comfort", "Will people feel comfortable?", AMBER),
    ("ASHRAE 62.1", "Fresh air", "Is there enough outside air to stay healthy?",
     GREEN),
    ("ASHRAE 90.1", "Energy", "Are we doing it without wasting energy?", ACCENT),
    ("Guideline 36", "The referee",
     "How should the equipment be run to hit all three?", INK),
]
y = Inches(2.15); rh = Inches(1.02); x0 = Inches(0.7); tot = Inches(11.93)
for i, (code, nick, q, c) in enumerate(rows):
    t = y + i * (rh + Inches(0.12))
    rect(s, x0, t, tot, rh, fill=CARD, line=LINE)
    rect(s, x0, t, Inches(2.5), rh, fill=c, line=None)
    box(s, x0 + Inches(0.2), t, Inches(2.1), rh, code, size=18, bold=True,
        color=WHITE, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    box(s, x0 + Inches(2.75), t, Inches(2.2), rh, nick, size=16, bold=True,
        color=c, anchor=MSO_ANCHOR.MIDDLE)
    box(s, x0 + Inches(5.0), t, tot - Inches(5.2), rh, q, size=15.5,
        color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
rect(s, Inches(0.7), Inches(6.75), Inches(11.93), Inches(0.6),
     fill=RGBColor(0xFF, 0xFB, 0xEB), line=RGBColor(0xFD, 0xE6, 0x8A))
box(s, Inches(0.98), Inches(6.82), Inches(11.4), Inches(0.45),
    "Comfort and fresh air push toward MORE conditioning; energy pushes "
    "toward LESS. Guideline 36 is the rulebook that settles the argument.",
    size=13.5, color=INK, anchor=MSO_ANCHOR.MIDDLE)

# =========================================================================
# 7 — ONE SHEET, FOUR OVERLAYS  (replaces the old Venn: the chart is not the
#     overlap of the standards, it is the sheet all of them are drawn on)
# =========================================================================
def mini_chart(s, X0, Y0, SX, SY):
    """Compact, unlabelled version of the chart for use as a motif."""
    px = lambda T: X0 + T * SX
    py = lambda W: Y0 - W * SY
    conn(s, px(0), py(0), px(40), py(0), color=AXIS, width=1.5)
    conn(s, px(0), py(0), px(0), py(30), color=AXIS, width=1.5)
    for (T0, W0), lw in [((13, 8.8), 0.9), ((33, 18.0), 0.9), ((24, 9.3), 1.3)]:
        sg = enth_seg(T0, W0)
        if sg:
            (t1, w1), (t2, w2) = sg
            conn(s, px(t1), py(w1), px(t2), py(w2), color=VIOLET, width=lw,
                 dash=True)
    v = [(Inches(px(t)), Inches(py(w))) for t, w in SAT]
    b = s.shapes.build_freeform(v[0][0], v[0][1], 1.0)
    b.add_line_segments(v[1:], close=False)
    cv = b.convert_to_shape()
    cv.fill.background(); cv.line.color.rgb = ACCENT; cv.line.width = Pt(2)
    cv.shadow.inherit = False
    cz = rect(s, Inches(px(20)), Inches(py(12)), Inches(px(26) - px(20)),
              Inches(py(6) - py(12)), fill=AMBER, line=AMBER, line_w=1.25)
    set_alpha(cz, 18)
    conn(s, px(33), py(18), px(24), py(9.3), color=GREEN, width=1.75, dash=True)
    conn(s, px(26.7), py(11.91), px(13), py(8.8), color=ACCENT, width=2)
    for (T, W), c in [((33, 18), AMBER), ((24, 9.3), GREEN),
                      ((26.7, 11.91), INK), ((13, 8.8), ACCENT)]:
        dot(s, px(T), py(W), c, r=0.05)


s = slide()
accent_header(s, "The big idea", "One sheet — every rulebook draws on it")

# hub: the chart itself
rect(s, Inches(4.75), Inches(2.1), Inches(3.85), Inches(4.0), fill=CARD,
     line=ACCENT, line_w=2.0)
box(s, Inches(4.85), Inches(2.3), Inches(3.65), Inches(0.35),
    "The psychrometric chart", size=16, bold=True, color=ACCENT,
    align=PP_ALIGN.CENTER)
box(s, Inches(4.85), Inches(2.66), Inches(3.65), Inches(0.3),
    "temperature × moisture — the only two axes", size=11, color=FAINT,
    align=PP_ALIGN.CENTER)
mini_chart(s, 5.15, 5.10, 0.072, 0.058)
box(s, Inches(4.9), Inches(5.32), Inches(3.55), Inches(0.6),
    "Every requirement around it becomes a shape on THESE SAME AXES.",
    size=11.5, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.2)

# four overlay cards
overlays = [
    (0.7, 2.1, AMBER, "ASHRAE 55 draws an AREA",
     "The comfort zone: the band of temperature and humidity occupants "
     "accept. The air-side of the standard becomes a region you must land "
     "the room inside.", None),
    (8.93, 2.1, GREEN, "ASHRAE 62.1 draws a POINT ON A LINE",
     "The required outside-air fraction fixes exactly where the mixed-air dot "
     "must sit along the line between return and outside air.", None),
    (0.7, 4.2, ACCENT, "ASHRAE 90.1 draws a DIRECTION",
     "The violet diagonals are equal-energy lines. Which side of them outside "
     "air falls on decides whether it is the cheaper source — and how far the "
     "dot may be pushed before energy is wasted.",
     "Air-side only — fan power and chiller efficiency sit outside the chart."),
    (8.93, 4.2, INK, "Guideline 36 draws the PATH",
     "The sequences are the route the dot travels through the day — "
     "economizer, coil, setpoint reset — from outside air to the room.", None),
]
for x, y, c, h, d, note in overlays:
    rect(s, Inches(x), Inches(y), Inches(3.7), Inches(1.9), fill=CARD,
         line=LINE)
    rect(s, Inches(x), Inches(y), Inches(3.7), Inches(0.1), fill=c, line=None,
         radius=False)
    box(s, Inches(x + 0.26), Inches(y + 0.22), Inches(3.2), Inches(0.32),
        h, size=13.5, bold=True, color=c)
    box(s, Inches(x + 0.26), Inches(y + 0.6), Inches(3.22), Inches(0.85),
        d, size=10.5, color=MUTED, line_spacing=1.15)
    if note:
        box(s, Inches(x + 0.26), Inches(y + 1.48), Inches(3.22), Inches(0.35),
            note, size=9.5, color=FAINT, line_spacing=1.1)

# arrows pointing in at the chart
for ax, adir in ((4.44, MSO_SHAPE.RIGHT_ARROW), (8.65, MSO_SHAPE.LEFT_ARROW)):
    for ay in (2.97, 5.07):
        ar = s.shapes.add_shape(adir, Inches(ax), Inches(ay), Inches(0.26),
                                Inches(0.16))
        ar.fill.solid(); ar.fill.fore_color.rgb = AXIS
        ar.line.fill.background(); ar.shadow.inherit = False

rect(s, Inches(0.7), Inches(6.3), Inches(11.93), Inches(0.95),
     fill=RGBColor(0xFF, 0xFB, 0xEB), line=RGBColor(0xFD, 0xE6, 0x8A))
box(s, Inches(0.98), Inches(6.4), Inches(11.4), Inches(0.3),
    "Read this the right way round", size=13, bold=True, color=AMBER)
box(s, Inches(0.98), Inches(6.69), Inches(11.5), Inches(0.5),
    "The chart is not the small overlap left where four standards happen to "
    "agree. It is the COMMON SHEET they are all drawn on — an area, a line, a "
    "direction and a path, sharing one pair of axes. Satisfy all four and the "
    "shapes simply fit together.", size=12.5, color=INK, line_spacing=1.15)

# =========================================================================
# 8 — ONE DOT ANSWERS ALL FOUR
# =========================================================================
s = slide()
accent_header(s, "How it satisfies each",
              "One dot on the chart answers every question")
items = [
    ("55 — Comfort",
     "The comfort zone is drawn straight onto the chart. Is the room dot "
     "inside the acceptable temperature/humidity band — and if not, which way "
     "is it off?", AMBER),
    ("62.1 — Fresh air",
     "The mixed-air dot must sit on the line between return and outside air. "
     "Its position along that line IS the outside-air proportion.", GREEN),
    ("90.1 — Energy",
     "Equal-energy lines reveal free cooling: when outside air holds less "
     "energy than return air, the economizer can meet the load with no "
     "mechanical cooling.", ACCENT),
    ("G36 — Sequences",
     "Visible proof the sequences land the dot inside the comfort zone, on "
     "the correct fresh-air line, along the lowest-energy path.", INK),
]
x0 = Inches(0.7); cw = Inches(5.85); ch = Inches(1.9)
gx = Inches(0.23); gy = Inches(0.22)
for i, (h, d, c) in enumerate(items):
    l = x0 + (i % 2) * (cw + gx); t = Inches(2.15) + (i // 2) * (ch + gy)
    rect(s, l, t, cw, ch, fill=CARD, line=LINE)
    rect(s, l, t, Inches(0.13), ch, fill=c, line=None, radius=False)
    box(s, l + Inches(0.32), t + Inches(0.2), cw - Inches(0.55), Inches(0.5),
        h, size=17, bold=True, color=c)
    box(s, l + Inches(0.32), t + Inches(0.72), cw - Inches(0.6),
        ch - Inches(0.85), d, size=13.5, color=MUTED, line_spacing=1.18)
box(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.6),
    "In one sentence: plot where the air IS, see where it MUST be, how it "
    "gets there, and whether it does so at least energy.",
    size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# =========================================================================
# 9 — COMPARISON (capability, not product)
# =========================================================================
s = slide()
accent_header(s, "The difference",
              "Trend-based BMS vs a psychrometric-aware BMS")
comp = [
    ("What the operator sees",
     "Several separate trend lines, read one at a time",
     "One picture showing the air's actual state"),
    ("Is the outside-air mix right?",
     "Inferred from the damper command",
     "The mixed-air dot must fall on the outside–return line"),
    ("Can we free-cool right now?",
     "Usually temperature-only, so humid-but-cool and dry-but-warm hours are "
     "misjudged",
     "Decided on true energy content of the air"),
    ("Humidity control",
     "A separate loop, often fighting the cooling loop",
     "Heat and moisture handled as one state"),
    ("Finding faults",
     "Alarm limits on individual points",
     "The shape of the process line exposes coil, damper and sensor faults"),
    ("Proving the standards",
     "Checked manually, usually after the fact",
     "Comfort, ventilation and energy judged together, live"),
]
tx = Inches(0.7); tw = Inches(11.93)
rect(s, tx, Inches(2.05), tw, Inches(4.75), fill=CARD, line=LINE)
c1 = Inches(0.3); c2 = Inches(3.4); c3 = Inches(8.15)
box(s, tx + c1, Inches(2.18), Inches(2.9), Inches(0.35),
    "QUESTION BEING ASKED", size=11, bold=True, color=FAINT)
box(s, tx + c2, Inches(2.18), Inches(4.5), Inches(0.35),
    "TREND-BASED BMS (TODAY)", size=11, bold=True, color=FAINT)
box(s, tx + c3, Inches(2.18), Inches(3.5), Inches(0.35),
    "PSYCHROMETRIC-AWARE BMS", size=11, bold=True, color=ACCENT)
ry = 2.62; rh = 0.68
for i, (k, a, b) in enumerate(comp):
    t = Inches(ry + i * rh)
    if i % 2 == 0:
        rect(s, tx + Inches(0.1), t, tw - Inches(0.2), Inches(rh),
             fill=RGBColor(0xFA, 0xFC, 0xFF), line=None, radius=False)
    box(s, tx + c1, t, Inches(3.0), Inches(rh), k, size=12.5, bold=True,
        color=INK, anchor=MSO_ANCHOR.MIDDLE)
    box(s, tx + c2, t, Inches(4.6), Inches(rh), a, size=12, color=MUTED,
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    box(s, tx + c3, t, Inches(3.5), Inches(rh), b, size=12, bold=True,
        color=GREEN, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
box(s, tx + Inches(0.3), Inches(6.15), Inches(11.3), Inches(0.5),
    "None of this needs new field equipment beyond humidity sensing on the "
    "air streams — the difference is what the software does with readings the "
    "building is already taking.", size=13, color=FAINT, line_spacing=1.15)

# =========================================================================
# 10 — CLOSING
# =========================================================================
s = slide(bg=RGBColor(0x0B, 0x12, 0x24))
rect(s, Inches(0), Inches(0), EMU_W, Inches(0.22), fill=ACCENT, line=None,
     radius=False)
box(s, Inches(1.0), Inches(2.4), Inches(11.3), Inches(1.2),
    "One chart. Four rulebooks.", size=42, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER)
box(s, Inches(1.4), Inches(3.9), Inches(10.5), Inches(1.5),
    "The psychrometric chart is not a design-only artifact. Bringing it into "
    "the building's controls turns anonymous setpoints back into visible "
    "physics — and puts comfort, air quality and energy into a single picture "
    "an operator can read at a glance.",
    size=18, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER,
    line_spacing=1.3)
box(s, Inches(2.0), Inches(5.6), Inches(9.3), Inches(1.0),
    "The computing power to draw it live has been ordinary for years. What "
    "has kept the chart out of the control room is habit and missing humidity "
    "sensors — not physics, and no longer technology.",
    size=15, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER,
    line_spacing=1.3)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Psychart-HVAC-ASHRAE-Overview.pptx")
prs.save(out)
print("Wrote", out, "-", len(prs.slides._sldIdLst), "slides")
